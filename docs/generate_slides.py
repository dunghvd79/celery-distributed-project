# -*- coding: utf-8 -*-
import sys
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    # Thiết lập kích thước slide 16:9 (Chuẩn Widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Định nghĩa bảng màu (Tech Palette - Tông sáng, hiện đại, tuyệt đối không dùng màu đen)
    c_navy = RGBColor(30, 58, 138)      # #1E3A8A (Royal Blue - Xanh công nghệ sáng và chuyên nghiệp)
    c_teal = RGBColor(13, 148, 136)    # #0D9488 (Teal - Celery Green)
    c_purple = RGBColor(124, 58, 237)  # #7C3AED (Bright Purple - Async)
    c_slate = RGBColor(71, 85, 105)     # #475569 (Slate text - Xám Slate sáng dịu mắt, không dùng màu đen)
    c_bg = RGBColor(255, 255, 255)     # #FFFFFF (Nền trắng tinh khiết - Sáng sủa)
    c_card_bg = RGBColor(248, 250, 252)# #F8FAFC (Xám nhạt Slate 50)
    c_border = RGBColor(226, 232, 240) # #E2E8F0 (Viền xám nhạt Slate 200)
    
    # Helper set background màu xám sáng
    def set_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = c_bg

    # Helper thêm tiêu đề tiêu chuẩn cho các slide nội dung
    def add_slide_header(slide, title_text, category="ỨNG DỤNG PHÂN TÁN (N05)"):
        # Category nhỏ ở trên
        tx_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = tx_cat.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = 'Arial'
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = c_teal
        
        # Tiêu đề chính
        tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = tx_title.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = c_navy

        # Accent Line mỏng màu teal nằm ngang dưới tiêu đề để tăng tính thẩm mỹ và hiện đại
        accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.03))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = c_teal
        accent_line.line.fill.background()

    # Helper vẽ một Card có bo góc và viền
    def draw_card(slide, left, top, width, height, title, content_lines, title_color=c_navy):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = c_card_bg
        shape.line.color.rgb = c_border
        shape.line.width = Pt(1.5)
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)
        
        # Tiêu đề của Card
        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = title_color
        p_title.space_after = Pt(10)
        
        # Nội dung của Card
        for line in content_lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.name = 'Arial'
            p.font.size = Pt(13)
            p.font.color.rgb = c_slate
            p.space_after = Pt(6)
            
            # Nếu bắt đầu bằng dấu gạch đầu dòng thì thụt vào tí
            if line.strip().startswith("-") or line.strip().startswith("•"):
                p.level = 0
        return shape

    # ==========================================
    # SLIDE 1: TRANG BÌA (Khớp với ảnh 1 của user)
    # ==========================================
    slide_layout = prs.slide_layouts[6] # Blank Layout
    slide1 = prs.slides.add_slide(slide_layout)
    set_bg(slide1)

    # Dải màu trang trí bên trái trang bìa (Sidebar Accent)
    bar1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = c_teal
    bar1.line.fill.background()
    
    bar2 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0), Inches(0.12), Inches(7.5))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = c_purple
    bar2.line.fill.background()
    
    # Title chính giữa phía trên
    tx_title1 = slide1.shapes.add_textbox(Inches(1.2), Inches(0.8), Inches(11.0), Inches(1.5))
    tf1 = tx_title1.text_frame
    tf1.word_wrap = True
    p1_title = tf1.paragraphs[0]
    p1_title.text = "Celery – Distributed Task Queue"
    p1_title.alignment = PP_ALIGN.CENTER
    p1_title.font.name = 'Arial'
    p1_title.font.size = Pt(44)
    p1_title.font.bold = True
    p1_title.font.color.rgb = c_navy
    
    p1_sub = tf1.add_paragraph()
    p1_sub.text = "Ứng dụng phân tán (N05) – Trường Công nghệ thông tin, Đại học Phenikaa"
    p1_sub.alignment = PP_ALIGN.CENTER
    p1_sub.font.name = 'Arial'
    p1_sub.font.size = Pt(16)
    p1_sub.font.bold = False
    p1_sub.font.color.rgb = c_slate
    p1_sub.space_before = Pt(12)

    # Khớp bố cục 3 Card của user (dịch chuyển nhẹ sang phải tránh sidebar)
    # Card 1: Giảng viên hướng dẫn
    draw_card(
        slide1, 
        left=Inches(1.6), 
        top=Inches(3.2), 
        width=Inches(4.9), 
        height=Inches(2.0),
        title="Giảng viên hướng dẫn",
        content_lines=["TS. Nguyễn Thành Trung"]
    )
    
    # Card 2: Thành viên Nhóm 2
    draw_card(
        slide1, 
        left=Inches(6.833), 
        top=Inches(3.2), 
        width=Inches(4.9), 
        height=Inches(2.0),
        title="Thành viên Nhóm 2",
        content_lines=[
            "Hoàng Văn Dũng – 23010438",
            "Nguyễn Hữu Thành Đạt – 23040102"
        ]
    )
    
    # Card 3: Thời gian (Chạy ngang bên dưới)
    draw_card(
        slide1, 
        left=Inches(1.6), 
        top=Inches(5.5), 
        width=Inches(10.133), 
        height=Inches(1.2),
        title="Thời gian thực hiện",
        content_lines=["Hà Nội, Tháng 5 năm 2026"]
    )

    # ==========================================
    # SLIDE 2: CELERY LÀ GÌ & TẠI SAO CẦN SỬ DỤNG
    # ==========================================
    slide2 = prs.slides.add_slide(slide_layout)
    set_bg(slide2)
    add_slide_header(slide2, "Celery là gì? Tại sao cần sử dụng?")
    
    draw_card(
        slide2,
        left=Inches(0.8),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        title="Celery là gì?",
        content_lines=[
            "• Là một hàng đợi tác vụ phân tán chạy nền (Distributed Task Queue) mã nguồn mở rất phổ biến trong hệ sinh thái Python.",
            "• Hoạt động dựa trên cơ chế truyền nhận thông điệp bất đồng bộ (Asynchronous Message Passing) giữa Client và Worker.",
            "• Hỗ trợ quản lý hàng đợi, lập lịch tác vụ (Periodic Tasks) và phân phối tải linh hoạt trên nhiều máy chủ khác nhau."
        ],
        title_color=c_teal
    )
    
    draw_card(
        slide2,
        left=Inches(6.933),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        title="Celery ra đời giải quyết vấn đề gì?",
        content_lines=[
            "• Giải phóng tài nguyên Web Server: Tránh hiện tượng block kết nối HTTP khi xử lý tác vụ nặng (như gửi email SMTP, nén ảnh). Web Server phản hồi ngay lập tức.",
            "• Ngăn nghẽn cổ chai: Tách biệt hoàn toàn tầng nhận yêu cầu (FastAPI/Django) và tầng thực thi nghiệp vụ (Celery Workers).",
            "• Tăng độ ổn định hệ thống: Cung cấp cơ chế tự động thử lại (Retry) khi các tác vụ gọi API bên thứ ba bị chập chờn."
        ],
        title_color=c_purple
    )

    # ==========================================
    # SLIDE 3: SƠ ĐỒ KIẾN TRÚC TỔNG THỂ (Khớp ảnh 2)
    # ==========================================
    slide3 = prs.slides.add_slide(slide_layout)
    set_bg(slide3)
    add_slide_header(slide3, "Sơ đồ kiến trúc tổng thể của Celery")
    
    # Vẽ các Tầng dạng Shape như sơ đồ của user
    # 1. Tầng sản xuất
    s_p = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(1.8), Inches(9.333), Inches(0.9))
    s_p.fill.solid()
    s_p.fill.fore_color.rgb = c_card_bg
    s_p.line.color.rgb = c_border
    s_p.line.width = Pt(1.5)
    s_p.text_frame.word_wrap = True
    p_sp_t = s_p.text_frame.paragraphs[0]
    p_sp_t.text = "TẦNG SẢN XUẤT (PRODUCERS)"
    p_sp_t.font.name = 'Arial'
    p_sp_t.font.bold = True
    p_sp_t.font.size = Pt(14)
    p_sp_t.font.color.rgb = c_teal
    p_sp_t.alignment = PP_ALIGN.CENTER
    p_sp_c = s_p.text_frame.add_paragraph()
    p_sp_c.text = "Các Web Application (FastAPI, Django, Flask) tạo và gửi yêu cầu tác vụ qua mạng"
    p_sp_c.font.name = 'Arial'
    p_sp_c.font.size = Pt(12)
    p_sp_c.font.color.rgb = c_slate
    p_sp_c.alignment = PP_ALIGN.CENTER
    
    # Label Arrow 1
    tx_a1 = slide3.shapes.add_textbox(Inches(2.0), Inches(2.72), Inches(9.333), Inches(0.4))
    tx_a1.text_frame.word_wrap = True
    p_a1 = tx_a1.text_frame.paragraphs[0]
    p_a1.text = "↓  1. Kích hoạt và Đóng gói Task (JSON)  ↓"
    p_a1.alignment = PP_ALIGN.CENTER
    p_a1.font.name = 'Arial'
    p_a1.font.size = Pt(11)
    p_a1.font.color.rgb = c_purple
    p_a1.font.italic = True
    p_a1.font.bold = True

    # 2. Tầng trung gian điều phối (Broker)
    s_b = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(3.12), Inches(9.333), Inches(1.0))
    s_b.fill.solid()
    s_b.fill.fore_color.rgb = c_card_bg
    s_b.line.color.rgb = c_border
    s_b.line.width = Pt(1.5)
    s_b.text_frame.word_wrap = True
    p_sb_t = s_b.text_frame.paragraphs[0]
    p_sb_t.text = "TẦNG TRUNG GIAN ĐIỀU PHỐI (MESSAGE BROKER)"
    p_sb_t.font.name = 'Arial'
    p_sb_t.font.bold = True
    p_sb_t.font.size = Pt(14)
    p_sb_t.font.color.rgb = c_teal
    p_sb_t.alignment = PP_ALIGN.CENTER
    p_sb_c = s_b.text_frame.add_paragraph()
    p_sb_c.text = "[ RabbitMQ / Redis Message Broker ]\n- Nhận thông điệp  - Xếp hàng đợi (Queue)  - Phân phối tin nhắn cho Worker rảnh"
    p_sb_c.font.name = 'Arial'
    p_sb_c.font.size = Pt(11)
    p_sb_c.font.color.rgb = c_slate
    p_sb_c.alignment = PP_ALIGN.CENTER

    # Label Arrow 2
    tx_a2 = slide3.shapes.add_textbox(Inches(2.0), Inches(4.14), Inches(9.333), Inches(0.4))
    tx_a2.text_frame.word_wrap = True
    p_a2 = tx_a2.text_frame.paragraphs[0]
    p_a2.text = "↓  2. Đẩy / Kéo Task theo cơ chế Fair-Dispatch  ↓"
    p_a2.alignment = PP_ALIGN.CENTER
    p_a2.font.name = 'Arial'
    p_a2.font.size = Pt(11)
    p_a2.font.color.rgb = c_purple
    p_a2.font.italic = True
    p_a2.font.bold = True

    # 3. Tầng xử lý phân tán (Workers)
    s_w = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(4.54), Inches(9.333), Inches(1.0))
    s_w.fill.solid()
    s_w.fill.fore_color.rgb = c_card_bg
    s_w.line.color.rgb = c_border
    s_w.line.width = Pt(1.5)
    s_w.text_frame.word_wrap = True
    p_sw_t = s_w.text_frame.paragraphs[0]
    p_sw_t.text = "TẦNG XỬ LÝ PHÂN TÁN (CELERY WORKERS)"
    p_sw_t.font.name = 'Arial'
    p_sw_t.font.bold = True
    p_sw_t.font.size = Pt(14)
    p_sw_t.font.color.rgb = c_teal
    p_sw_t.alignment = PP_ALIGN.CENTER
    p_sw_c = s_w.text_frame.add_paragraph()
    p_sw_c.text = "Các Worker Node chạy song song trên nhiều máy chủ độc lập\n[ Worker Node 1 ]       [ Worker Node 2 ]       [ Worker Node N ]"
    p_sw_c.font.name = 'Arial'
    p_sw_c.font.size = Pt(11)
    p_sw_c.font.color.rgb = c_slate
    p_sw_c.alignment = PP_ALIGN.CENTER

    # Label Arrow 3
    tx_a3 = slide3.shapes.add_textbox(Inches(2.0), Inches(5.56), Inches(9.333), Inches(0.4))
    tx_a3.text_frame.word_wrap = True
    p_a3 = tx_a3.text_frame.paragraphs[0]
    p_a3.text = "↓  3. Ghi nhận trạng thái xử lý và kết quả  ↓"
    p_a3.alignment = PP_ALIGN.CENTER
    p_a3.font.name = 'Arial'
    p_a3.font.size = Pt(11)
    p_a3.font.color.rgb = c_purple
    p_a3.font.italic = True
    p_a3.font.bold = True

    # 4. Tầng lưu trữ kết quả (Backend)
    s_r = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(5.96), Inches(9.333), Inches(0.8))
    s_r.fill.solid()
    s_r.fill.fore_color.rgb = c_card_bg
    s_r.line.color.rgb = c_border
    s_r.line.width = Pt(1.5)
    s_r.text_frame.word_wrap = True
    p_sr_t = s_r.text_frame.paragraphs[0]
    p_sr_t.text = "TẦNG LƯU TRỮ KẾT QUẢ (RESULT BACKEND)"
    p_sr_t.font.name = 'Arial'
    p_sr_t.font.bold = True
    p_sr_t.font.size = Pt(14)
    p_sr_t.font.color.rgb = c_teal
    p_sr_t.alignment = PP_ALIGN.CENTER
    p_sr_c = s_r.text_frame.add_paragraph()
    p_sr_c.text = "[ Redis Database / PostgreSQL / Memcached ] để lưu giá trị trả về và thông tin lỗi"
    p_sr_c.font.name = 'Arial'
    p_sr_c.font.size = Pt(11)
    p_sr_c.font.color.rgb = c_slate
    p_sr_c.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 4: CORE FEATURES (1/2) - ASYNC & CPU-BOUND
    # ==========================================
    slide4 = prs.slides.add_slide(slide_layout)
    set_bg(slide4)
    add_slide_header(slide4, "Các tính năng cốt lõi của Celery (Phần 1)")
    
    draw_card(
        slide4,
        left=Inches(0.8),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        title="1. Xử lý tác vụ Bất đồng bộ (Async - TN1)",
        content_lines=[
            "• Luồng hoạt động:",
            "  FastAPI nhận HTTP request -> Đẩy tác vụ vào Broker -> Trả ngay mã Task ID cho Client (Response Time < 10ms).",
            "• Ví dụ thực tế: Tác vụ gửi Email hàng loạt.",
            "  Giúp giải phóng tài nguyên Web Server tức thời, loại bỏ thời gian chờ đợi phản hồi từ máy chủ SMTP.",
            "• Thống kê hiệu năng:",
            "  Thời gian phản hồi API giảm mạnh từ ~3.0 giây xuống dưới 0.01 giây."
        ],
        title_color=c_teal
    )
    
    draw_card(
        slide4,
        left=Inches(6.933),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        title="2. Xử lý tác vụ nặng độc lập (CPU-bound - TN2)",
        content_lines=[
            "• Luồng hoạt động:",
            "  Tác vụ xử lý ảnh dung lượng lớn, tạo ảnh đại diện (thumbnail) được ủy thác cho Worker xử lý độc lập.",
            "• Tận dụng tài nguyên:",
            "  Chạy ngầm tách biệt hoàn toàn giúp Web Server chính không bị quá tải tài nguyên hệ thống hoặc hết luồng xử lý CPU.",
            "• Trải nghiệm người dùng:",
            "  Giao diện Web hoạt động mượt mà, người dùng có thể tiếp tục thao tác các luồng chức năng khác trên ứng dụng."
        ],
        title_color=c_purple
    )

    # ==========================================
    # SLIDE 5: CORE FEATURES (2/2) - PROGRESS & WORKFLOWS
    # ==========================================
    slide5 = prs.slides.add_slide(slide_layout)
    set_bg(slide5)
    add_slide_header(slide5, "Các tính năng cốt lõi của Celery (Phần 2)")
    
    draw_card(
        slide5,
        left=Inches(0.8),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        title="3. Theo dõi tiến trình thời gian thực (Progress - TN3)",
        content_lines=[
            "• Luồng hoạt động:",
            "  Worker cập nhật trạng thái tác vụ lên Redis thông qua API: self.update_state(state='PROGRESS', meta=...).",
            "• Cập nhật động:",
            "  Giúp client liên tục gửi request truy vấn tiến độ (10%, 40%, 70%, 90%) dựa trên Task ID.",
            "• Ứng dụng thực tế:",
            "  Vẽ thanh tiến trình trực quan cho các tác vụ tốn nhiều thời gian như truy vấn cơ sở dữ liệu lớn và xuất file báo cáo PDF."
        ],
        title_color=c_teal
    )
    
    draw_card(
        slide5,
        left=Inches(6.933),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        title="4. Định nghĩa quy trình phối hợp (Workflows - TN4)",
        content_lines=[
            "• Celery Group (TN4.1 - Song song):",
            "  Nhóm nhiều tác vụ chạy song song cùng lúc (ví dụ: tạo 3 kích cỡ ảnh avatar, banner, thumbnail). Tối ưu hóa tối đa thời gian xử lý tổng thể.",
            "• Celery Chain (TN4.2 - Tuần tự):",
            "  Kết nối các tác vụ chạy nối tiếp nhau, kết quả đầu ra của tác vụ trước là đầu vào của tác vụ sau.",
            "• Ý nghĩa:",
            "  Xây dựng các luồng xử lý phức tạp một cách dễ dàng thông qua cú pháp code Python ngắn gọn."
        ],
        title_color=c_purple
    )

    # ==========================================
    # SLIDE 6: KHÓA PHÂN TÁN (TN5) - ĐẶT VẤN ĐỀ
    # ==========================================
    slide6 = prs.slides.add_slide(slide_layout)
    set_bg(slide6)
    add_slide_header(slide6, "Tính năng nâng cao: Khóa phân tán (Distributed Lock)")
    
    draw_card(
        slide6,
        left=Inches(0.8),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        title="Đặt vấn đề: Bài toán Race Condition",
        content_lines=[
            "• Tình huống xảy ra:",
            "  Khi người dùng nhấn nút 'Thanh toán' liên tiếp 2 lần do mạng lag hoặc cố tình gửi đồng thời các yêu cầu giao dịch.",
            "• Hậu quả thực tế:",
            "  Hệ thống xử lý hai tác vụ cùng lúc mà không có cơ chế chặn, dẫn đến việc tài khoản người dùng bị trừ tiền 2 lần (Race Condition).",
            "• Yêu cầu bài toán:",
            "  Phải thiết lập cơ chế khóa đồng bộ trên toàn mạng lưới phân tán để chỉ cho phép một tác vụ chạy tại một thời điểm."
        ],
        title_color=c_navy
    )
    
    draw_card(
        slide6,
        left=Inches(6.933),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        title="Hạn chế của cơ chế khóa cục bộ (Local Lock)",
        content_lines=[
            "• Không hoạt động trong môi trường phân tán:",
            "  Các cơ chế khóa đồng bộ thông thường (như threading.Lock trong Python) chỉ có tác dụng trong phạm vi 1 tiến trình.",
            "• Thách thức với Multi-Worker:",
            "  Khi các tác vụ chạy trên nhiều máy chủ Worker độc lập khác nhau, chúng không thể giao tiếp trực tiếp để khóa lẫn nhau.",
            "• Giải pháp:",
            "  Cần một hệ thống quản lý khóa tập trung có tốc độ truy xuất cực nhanh nằm ngoài phạm vi các Worker (Redis Lock)."
        ],
        title_color=c_teal
    )

    # ==========================================
    # SLIDE 7: KHÓA PHÂN TÁN (TN5) - GIẢI PHÁP & KẾT QUẢ
    # ==========================================
    slide7 = prs.slides.add_slide(slide_layout)
    set_bg(slide7)
    add_slide_header(slide7, "Giải pháp Khóa phân tán sử dụng Redis Lock")
    
    draw_card(
        slide7,
        left=Inches(0.8),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        title="Cơ chế giải quyết của Redis Lock (TN5)",
        content_lines=[
            "1. Tác vụ 1 chạy -> Thiết lập một khóa duy nhất trên Redis gắn với mã hóa đơn (ví dụ: lock:payment:TXN_12345).",
            "2. Tác vụ 2 chạy song song -> Lên Redis kiểm tra thấy khóa đã tồn tại.",
            "3. Tác vụ 2 chủ động hủy bỏ thực thi ngay lập tức thông qua câu lệnh: raise Ignore().",
            "4. Tác vụ 1 hoàn thành thanh toán -> Giải phóng khóa trên Redis để đón nhận giao dịch mới."
        ],
        title_color=c_teal
    )
    
    draw_card(
        slide7,
        left=Inches(6.933),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        title="Kết quả thử nghiệm thực tế",
        content_lines=[
            "• Bảo vệ tài khoản an toàn 100%:",
            "  Tác vụ trùng lặp bị chặn đứng ngay ở tầng kiểm tra khóa, không thực hiện hành vi trừ tiền.",
            "• Tối ưu trạng thái Celery (IGNORED):",
            "  Việc sử dụng trạng thái IGNORED giúp phân biệt rõ ràng giữa 'Tác vụ lỗi' (Failure) và 'Tác vụ bị chủ động hủy bỏ để bảo vệ' (Ignored).",
            "• Hiệu suất cao:",
            "  Thời gian kiểm tra khóa phân tán trên RAM của Redis chỉ mất chưa đầy 1 mili-giây."
        ],
        title_color=c_purple
    )

    # ==========================================
    # SLIDE 8: DEAD LETTER QUEUE (TN6) - ĐẶT VẤN ĐỀ
    # ==========================================
    slide8 = prs.slides.add_slide(slide_layout)
    set_bg(slide8)
    add_slide_header(slide8, "Tính năng nâng cao: Hàng đợi thư chết (DLQ)")
    
    draw_card(
        slide8,
        left=Inches(0.8),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        title="Đặt vấn đề: Sự không ổn định của APIs bên thứ ba",
        content_lines=[
            "• Tình huống xảy ra:",
            "  Worker xử lý tác vụ đồng bộ thông tin đơn hàng hoặc gọi API thanh toán của bên thứ ba, nhưng đối tác bị mất kết nối mạng hoặc sập máy chủ tạm thời.",
            "• Rủi ro mất mát dữ liệu:",
            "  Nếu tác vụ lập tức bị coi là thất bại, dữ liệu đơn hàng sẽ bị mất hoặc để lại trạng thái lỗi hệ thống.",
            "• Yêu cầu giải pháp:",
            "  Worker cần tự động thử lại (Retry). Nếu thử lại hết lượt mà vẫn lỗi, phải lưu trữ thông tin lỗi an toàn để xử lý lại."
        ],
        title_color=c_navy
    )
    
    draw_card(
        slide8,
        left=Inches(6.933),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        title="Ý niệm về Hàng đợi thư chết (Dead Letter Queue)",
        content_lines=[
            "• Hàng đợi độc lập (`celery.dlq`):",
            "  Là một hàng đợi riêng biệt trong RabbitMQ, chuyên dùng để chứa các gói tin tác vụ đã thất bại hoàn toàn sau nhiều lần thử lại.",
            "• Bảo toàn dữ liệu lỗi:",
            "  Gói tin chứa tham số ban đầu được giữ nguyên vẹn trong hàng đợi DLQ, không bị vứt bỏ, giúp tránh mất dữ liệu khách hàng.",
            "• Phân tách lỗi:",
            "  Giúp hàng đợi chính luôn thông suốt, không bị nghẽn bởi các tin nhắn liên tục bị lỗi kẹt lại."
        ],
        title_color=c_teal
    )

    # ==========================================
    # SLIDE 9: DEAD LETTER QUEUE (TN6) - GIẢI PHÁP & KẾT QUẢ
    # ==========================================
    slide9 = prs.slides.add_slide(slide_layout)
    set_bg(slide9)
    add_slide_header(slide9, "Giải pháp Tự động Retry & Cô lập lỗi vào DLQ")
    
    draw_card(
        slide9,
        left=Inches(0.8),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        title="Cơ chế tự động xử lý của hệ thống (TN6)",
        content_lines=[
            "1. Tác vụ Flaky lỗi mạng -> Tự động thử lại (max_retries=3, tổng cộng 4 lần chạy thử).",
            "2. Đếm ngược thử lại: Mỗi lần cách nhau 5 giây để đợi mạng hoặc API đối tác phục hồi.",
            "3. Khi đạt giới hạn 4 lần vẫn lỗi -> Hệ thống tự động kích hoạt đẩy thông tin lỗi vào hàng đợi celery.dlq.",
            "4. Bắn thông báo tự động (Slack/Telegram) báo cáo sự cố kèm mã lỗi cho Quản trị viên."
        ],
        title_color=c_teal
    )
    
    draw_card(
        slide9,
        left=Inches(6.933),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        title="Kết quả & Tính thực tiễn",
        content_lines=[
            "• Bảo toàn tuyệt đối dữ liệu:",
            "  Tin nhắn lỗi được cô lập an toàn trong RabbitMQ (chứa tham số thô dạng JSON).",
            "• Giải phóng và xử lý lại (Ack / Replay):",
            "  Khi API đối tác sửa xong, Quản trị viên chỉ cần ấn nút 'Giải phóng DLQ' để lấy toàn bộ các đơn kẹt ra xử lý lại hàng loạt cùng lúc.",
            "• Giám sát trực quan:",
            "  Dashboard và RabbitMQ hiển thị chuẩn xác số lượng tin nhắn lỗi đang bị cô lập trong DLQ."
        ],
        title_color=c_purple
    )

    # ==========================================
    # SLIDE 10: TỔNG KẾT & HƯỚNG PHÁT TRIỂN
    # ==========================================
    slide10 = prs.slides.add_slide(slide_layout)
    set_bg(slide10)
    add_slide_header(slide10, "Tổng kết & Hướng phát triển dự án")
    
    draw_card(
        slide10,
        left=Inches(0.8),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        title="Kết luận thực nghiệm",
        content_lines=[
            "• Bất đồng bộ hóa thành công: FastAPI phản hồi tức thời (<0.01s), Celery Worker tự động xử lý ngầm.",
            "• Kiểm soát tranh chấp an toàn: Khóa phân tán Redis triệt tiêu hoàn toàn lỗi Race Condition thanh toán trùng lặp.",
            "• Đảm bảo độ bền bỉ dữ liệu: Hàng đợi thư chết (DLQ) và Slack Alert giúp cách ly và xử lý sự cố API thông minh.",
            "• Thiết kế tối giản: Dashboard tiếng Việt hỗ trợ kéo dãn bảng logs hệ thống tiện lợi."
        ],
        title_color=c_teal
    )
    
    draw_card(
        slide10,
        left=Inches(6.933),
        top=Inches(1.8),
        width=Inches(5.6),
        height=Inches(4.8),
        title="Hướng phát triển tương lai",
        content_lines=[
            "• Auto-scaling Workers:",
            "  Tự động scale tăng/giảm số lượng tiến trình Worker dựa trên lưu lượng hàng đợi tin nhắn trong RabbitMQ.",
            "• Tích hợp Prometheus & Grafana:",
            "  Giám sát trực quan hóa hiệu năng của các Worker core thời gian thực thay vì sử dụng Flower cơ bản.",
            "• Rate Limiting nâng cao:",
            "  Cấu hình giới hạn tần suất gọi API bên thứ ba tự động để tránh bị chặn IP."
        ],
        title_color=c_purple
    )

    # Tạo thư mục docs nếu chưa tồn tại
    if not os.path.exists("docs"):
        os.makedirs("docs")
        
    # Lưu tệp tin trình bày
    output_path = "docs/project_presentation.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
